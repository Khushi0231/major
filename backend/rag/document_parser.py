"""Complete document parser supporting PDF, DOCX, PPTX, TXT, images, and code files"""
import os
import logging
from typing import List, Dict, Tuple
from pathlib import Path

logger = logging.getLogger(__name__)

# PDF parsing
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    try:
        import pypdf
        PYMUPDF_AVAILABLE = False
        PYPDF_AVAILABLE = True
    except ImportError:
        PYMUPDF_AVAILABLE = False
        PYPDF_AVAILABLE = False
        logger.warning("PyMuPDF and pypdf not available")

try:
    from docling.datamodel.base_models import InputFormat
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    logger.warning("Docling not available, will use PyMuPDF fallback")

# DOCX parsing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX parsing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Image OCR
try:
    from PIL import Image
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    logger.warning("Tesseract OCR not available")


def extract_pdf(path: str) -> List[Dict[str, any]]:
    """Extract text from PDF using Docling (fallback to PyMuPDF)"""
    pages = []
    
    # Try Docling first
    if DOCLING_AVAILABLE:
        try:
            converter = DocumentConverter()
            result = converter.convert(path)
            doc_text = result.document.export_to_text()
            # Split by pages if possible
            lines = doc_text.split('\n')
            chunk_size = 50  # Approximate lines per page
            for i in range(0, len(lines), chunk_size):
                page_text = '\n'.join(lines[i:i+chunk_size])
                if page_text.strip():
                    pages.append({"page": (i // chunk_size) + 1, "text": page_text})
            if pages:
                logger.info(f"Extracted {len(pages)} pages from PDF using Docling")
                return pages
        except Exception as e:
            logger.warning(f"Docling failed, falling back to PyMuPDF: {e}")
    
    # Fallback to PyMuPDF
    if PYMUPDF_AVAILABLE:
        try:
            doc = fitz.open(path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    pages.append({"page": page_num + 1, "text": text})
            doc.close()
            logger.info(f"Extracted {len(pages)} pages from PDF using PyMuPDF")
            return pages
        except Exception as e:
            logger.error(f"PyMuPDF extraction failed: {e}")
    
    # Fallback to pypdf
    if PYPDF_AVAILABLE:
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    pages.append({"page": page_num, "text": text})
            logger.info(f"Extracted {len(pages)} pages from PDF using pypdf")
            return pages
        except Exception as e:
            logger.error(f"pypdf extraction failed: {e}")
    
    raise Exception("PDF extraction failed: No available PDF parser")


def extract_docx(path: str) -> List[Dict[str, any]]:
    """Extract text from DOCX file"""
    if not DOCX_AVAILABLE:
        raise Exception("python-docx not available")
    
    try:
        doc = Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        text = "\n".join(paragraphs)
        return [{"page": 1, "text": text}]
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        raise


def extract_pptx(path: str) -> List[Dict[str, any]]:
    """Extract text from PPTX file"""
    if not PPTX_AVAILABLE:
        raise Exception("python-pptx not available")
    
    try:
        prs = Presentation(path)
        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                slides.append({"page": slide_num, "text": "\n".join(slide_text)})
        return slides
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        raise


def extract_image(path: str) -> List[Dict[str, any]]:
    """Extract text from image using OCR"""
    if not TESSERACT_AVAILABLE:
        raise Exception("Tesseract OCR not available")
    
    try:
        image = Image.open(path)
        text = pytesseract.image_to_string(image)
        if text.strip():
            return [{"page": 1, "text": text}]
        return []
    except Exception as e:
        logger.error(f"Image OCR failed: {e}")
        raise


def extract_code_file(path: str) -> List[Dict[str, any]]:
    """Extract text from code files"""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return [{"page": 1, "text": content}]
    except Exception as e:
        logger.error(f"Code file extraction failed: {e}")
        raise


def extract_text_file(path: str) -> List[Dict[str, any]]:
    """Extract text from plain text files"""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return [{"page": 1, "text": content}]
    except Exception as e:
        logger.error(f"Text file extraction failed: {e}")
        raise


def parse_document(file_path: str) -> List[Dict[str, any]]:
    """
    Main document parser that routes to appropriate extractor based on file extension.
    Returns list of {"page": int, "text": str} dictionaries.
    """
    ext = Path(file_path).suffix.lower()
    
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp"]:
        return extract_image(file_path)
    elif ext in [".py", ".java", ".cpp", ".js", ".json"]:
        return extract_code_file(file_path)
    elif ext in [".txt", ".md"]:
        return extract_text_file(file_path)
    else:
        raise Exception(f"Unsupported file type: {ext}")


def chunk_text_for_storage(pages: List[Dict[str, any]], chunk_size: int = 512, overlap: int = 50) -> List[Tuple[str, Dict]]:
    """
    Chunk text from pages into smaller pieces for vector storage.
    Returns list of (chunk_text, metadata_dict) tuples.
    """
    out_chunks = []
    
    for page_data in pages:
        text = page_data.get("text", "")
        page_num = page_data.get("page", 1)
        
        # Split text into chunks
        words = text.split()
        if not words:
            continue
        
        start = 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk_words = words[start:end]
            chunk_text = " ".join(chunk_words).strip()
            
            if chunk_text:
                metadata = {
                    "page": page_num,
                    "chunk_start": start,
                    "chunk_end": end
                }
                out_chunks.append((chunk_text, metadata))
            
            # Move forward with overlap
            start = max(end - overlap, start + 1) if end - overlap > start else end
    
    return out_chunks
