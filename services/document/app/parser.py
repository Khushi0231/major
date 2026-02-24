"""Document Service - File parser (PDF, DOCX, PPTX, TXT, images)"""
import logging
from typing import List, Dict, Tuple
from pathlib import Path

logger = logging.getLogger(__name__)

# ─── Optional deps with graceful fallback ──────
try:
    import fitz  # PyMuPDF
    PYMUPDF = True
except ImportError:
    PYMUPDF = False

try:
    from docx import Document
    DOCX = True
except ImportError:
    DOCX = False

try:
    from pptx import Presentation
    PPTX = True
except ImportError:
    PPTX = False

try:
    from PIL import Image
    import pytesseract
    TESSERACT = True
except ImportError:
    TESSERACT = False


def parse_document(file_path: str) -> List[Dict]:
    """Route to correct extractor. Returns [{page, text}, ...]."""
    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".jpg": _extract_image,
        ".jpeg": _extract_image,
        ".png": _extract_image,
    }
    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return fn(file_path)


def chunk_text(pages: List[Dict], chunk_size: int = 500, overlap: int = 50) -> List[Tuple[str, Dict]]:
    """Split page text into overlapping chunks. Returns [(text, metadata), ...]."""
    out = []
    for page in pages:
        words = page.get("text", "").split()
        if not words:
            continue
        start = 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk = " ".join(words[start:end]).strip()
            if chunk:
                out.append((chunk, {"page": page.get("page", 1), "start": start, "end": end}))
            start = max(end - overlap, start + 1) if end - overlap > start else end
    return out


# ─── Extractors ────────────────────────────────

def _extract_pdf(path: str) -> List[Dict]:
    if not PYMUPDF:
        raise RuntimeError("PyMuPDF not installed")
    doc = fitz.open(path)
    pages = []
    for i in range(len(doc)):
        text = doc[i].get_text()
        if text.strip():
            pages.append({"page": i + 1, "text": text})
    doc.close()
    return pages


def _extract_docx(path: str) -> List[Dict]:
    if not DOCX:
        raise RuntimeError("python-docx not installed")
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"page": 1, "text": text}] if text else []


def _extract_pptx(path: str) -> List[Dict]:
    if not PPTX:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if parts:
            slides.append({"page": i, "text": "\n".join(parts)})
    return slides


def _extract_text(path: str) -> List[Dict]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    return [{"page": 1, "text": content}] if content.strip() else []


def _extract_image(path: str) -> List[Dict]:
    if not TESSERACT:
        raise RuntimeError("pytesseract not installed")
    text = pytesseract.image_to_string(Image.open(path))
    return [{"page": 1, "text": text}] if text.strip() else []
